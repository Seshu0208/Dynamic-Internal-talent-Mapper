import pandas as pd
import docx
from pptx import Presentation
import pdfplumber
import PyPDF2
import os

def read_txt(path):
    with open(path, "r", encoding="utf-8") as f:
        return f.read()

def read_csv(path):
    df = pd.read_csv(path)
    return df.to_string()

def read_excel(path):
    df = pd.read_excel(path)
    return df.to_string()

def read_docx(path):
    doc = docx.Document(path)
    return "\n".join([p.text for p in doc.paragraphs])

def read_pptx(path):
    prs = Presentation(path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def read_pdf(path):
    text = ""
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            text += page.extract_text() or ""
    return text

def extract_text_from_file(path):
    ext = os.path.splitext(path)[1].lower()

    if ext == ".txt":
        return read_txt(path)
    elif ext == ".csv":
        return read_csv(path)
    elif ext in [".xls", ".xlsx"]:
        return read_excel(path)
    elif ext == ".docx":
        return read_docx(path)
    elif ext == ".pptx":
        return read_pptx(path)
    elif ext == ".pdf":
        return read_pdf(path)
    else:
        raise ValueError("Unsupported file format")